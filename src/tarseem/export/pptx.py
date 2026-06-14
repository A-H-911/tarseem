"""PPTX writer — native PowerPoint shapes/connectors from the positioned IR (08 §3, D2:C).

Consumes the positioned IR (ADR-001: writers never lay out) and emits **native** python-pptx
autoshapes, connectors, freeform edges, and text — never an embedded image or SVG-ungroup
(invariant 5). Geometry is the IR's pixel coordinates converted to EMU (9525 EMU/px @ 96 dpi).
The SVG framing is matched per family: generic/ER diagrams translate by a 24px margin, swimlane/
sequence bake absolute coordinates, so the slide margin follows the same rule.

Constraints honoured:
- **RTL** (invariant 4): python-pptx has no API for paragraph direction or the complex-script
  font, so RTL labels get ``<a:pPr rtl="1">`` and every run's ``a:cs``/``a:ea`` typeface is set
  to Cairo (otherwise Arabic falls back to the theme's complex-script font) via lxml patches.
- **Determinism** (invariant 7): a .pptx is a zip that python-pptx stamps with wall-clock mtimes
  + core-property timestamps; core props are pinned to a constant and the zip is re-emitted with
  normalized ``ZipInfo``, so the same spec ⇒ byte-identical .pptx.
- **Capability reports, never silent drops** (invariant 6): see ``_report``.
"""
from __future__ import annotations

import io
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from tarseem.export.result import WriteResult
from tarseem.model.ir import LaneBand, Marker, PhaseBand, PositionedDiagram, PositionedEdge
from tarseem.model.ir import PositionedNode as PNode
from tarseem.render.text import (
    has_rtl,
    resolve_badge_side,
    resolve_direction,
    resolve_edge_corners,
)
from tarseem.report import CapabilityWarning, build_capability_report

if TYPE_CHECKING:  # `Presentation` is a factory fn; the instance type lives here
    from pptx.presentation import Presentation as _Prs

__all__ = ["write_pptx", "to_pptx_bytes"]

EMU_PER_PX = 9525  # 914400 EMU/inch ÷ 96 px/inch
_FONT = "Cairo"  # names the SVG face; PowerPoint substitutes if absent (fonts ceiling)
_FIXED_TS = datetime(2001, 1, 1, tzinfo=timezone.utc)  # constant (invariant 7: no wall-clock)
_EDGE_RADIUS = 8.0  # corner-rounding radius — MUST match render/svg.py _EDGE_RADIUS

# Colours — MUST match the SVG/draw.io writers so all writers agree.
_DEFAULT_FILL = "#FFFFFF"
_DEFAULT_STROKE = "#333333"
_DEFAULT_TEXT = "#14281D"
_DEFAULT_EDGE = "#333333"
_SEP = "#B0BEC5"
_PHASE_FILL = "#37474F"
_TITLE_FILL = "#269973"
_MARKER_BLACK = "#000000"
_LANE_ROW = "#EEEEEE"
_LANE_ACCENT = "#333333"
_ER_TITLE_FILL = "#37474F"
_ER_BORDER = "#5A6B7B"
_ER_ROW_SEP = "#CFD8DC"
_ER_KEY_FILL = {"PK": "#C49000", "FK": "#3B7DD8"}
_SEQ_STEM = "#9AA8A2"
_SEQ_ACT_BORDER = "#2E8B57"

_BADGE_R = 11.0
_CUBE_DEPTH = 14.0
_LABEL_W = 160.0
_V_HEADER = 64.0
_CHIP_H = 56.0
_V_CHIP_H = 48.0
_CHIP_INSET = 8.0
_EDGE_WIDTH_DEFAULT = {"swimlane": 2.0, "er": 1.5, "sequence": 1.5}
_SHADOW_SHAPES = ("cube", "cylinder")  # 3-D shapes get a drop shadow (matches SVG/draw.io)

# python-pptx autoshape per IR shape (documented MSO_SHAPE members only).
_SHAPE: dict[str, MSO_SHAPE] = {
    "rect": MSO_SHAPE.RECTANGLE,
    "roundrect": MSO_SHAPE.ROUNDED_RECTANGLE,
    "rounded": MSO_SHAPE.ROUNDED_RECTANGLE,
    "stadium": MSO_SHAPE.ROUNDED_RECTANGLE,  # adjustment maxed -> stadium
    "diamond": MSO_SHAPE.DIAMOND,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "cylinder": MSO_SHAPE.CAN,
    "document": MSO_SHAPE.FLOWCHART_DOCUMENT,
    "cube": MSO_SHAPE.CUBE,
}


def _rgb(hexcolor: str) -> RGBColor:
    return RGBColor.from_string(str(hexcolor).lstrip("#")[:6].upper())


def _fill_of(style: dict) -> str:
    return str(style.get("fill", _DEFAULT_FILL))


def _stroke_of(style: dict) -> str:
    return str((style.get("border") or {}).get("color", _DEFAULT_STROKE))


def _stroke_w_of(style: dict) -> float:
    return float((style.get("border") or {}).get("width", 1) or 1)


def _font_color_of(style: dict) -> str:
    return str((style.get("text") or {}).get("color") or _DEFAULT_TEXT)


def _rtl_label(label) -> bool:
    if label is None:
        return False
    return resolve_direction(label.direction, label.text) == "rtl" or has_rtl(label.text)


def _toward(b: tuple[float, float], a: tuple[float, float], r: float) -> tuple[float, float]:
    """Point ``r`` from ``b`` toward ``a`` (clamped to half the segment). MUST match svg."""
    dx, dy = a[0] - b[0], a[1] - b[1]
    dist = (dx * dx + dy * dy) ** 0.5 or 1.0
    rr = min(r, dist / 2)
    return (b[0] + dx / dist * rr, b[1] + dy / dist * rr)


def _rounded_points(points: list[tuple[float, float]], curved: bool) -> list[tuple[float, float]]:
    """Densify a polyline with rounded corners so a straight-segment freeform looks curved like
    the SVG (which draws a quadratic at each bend). Each corner's quadratic is sampled into short
    segments. ``curved=False`` (theme.edgeCorners=straight) returns the raw points."""
    if not curved or len(points) <= 2:
        return list(points)
    out: list[tuple[float, float]] = [points[0]]
    for i in range(1, len(points) - 1):
        a, b, c = points[i - 1], points[i], points[i + 1]
        p0, p2 = _toward(b, a, _EDGE_RADIUS), _toward(b, c, _EDGE_RADIUS)
        out.append(p0)
        steps = 6
        for k in range(1, steps + 1):
            t = k / steps
            mt = 1 - t
            out.append((
                mt * mt * p0[0] + 2 * mt * t * b[0] + t * t * p2[0],
                mt * mt * p0[1] + 2 * mt * t * b[1] + t * t * p2[1],
            ))
    out.append(points[-1])
    return out


def _set_run_font(run, name: str) -> None:
    """Set latin + east-asian + complex-script typefaces. python-pptx's ``font.name`` only sets
    ``a:latin``; Arabic renders from ``a:cs``, so without this it falls back to the theme font."""
    run.font.name = name  # a:latin, inserted in the schema-correct position
    rPr = run.font._rPr
    latin = rPr.find(qn("a:latin"))
    for tag in ("a:ea", "a:cs"):
        for existing in rPr.findall(qn(tag)):
            rPr.remove(existing)
    latin.addnext(rPr.makeelement(qn("a:cs"), {"typeface": name}))  # -> latin, cs
    latin.addnext(rPr.makeelement(qn("a:ea"), {"typeface": name}))  # -> latin, ea, cs


def _add_shadow(sp) -> None:
    """A subtle outer drop shadow (owner-preferred for 3-D shapes; mirrored in SVG/draw.io)."""
    spPr = sp._element.spPr
    for existing in spPr.findall(qn("a:effectLst")):
        spPr.remove(existing)
    eff = spPr.makeelement(qn("a:effectLst"), {})
    shdw = eff.makeelement(
        qn("a:outerShdw"),
        {"blurRad": "40000", "dist": "28000", "dir": "5400000", "rotWithShape": "0"},
    )
    clr = shdw.makeelement(qn("a:srgbClr"), {"val": "000000"})
    clr.append(clr.makeelement(qn("a:alpha"), {"val": "32000"}))
    shdw.append(clr)
    eff.append(shdw)
    spPr.append(eff)


def _set_alpha(fore_color, pct: int) -> None:
    """Set fill alpha (0–100) via the srgbClr's a:alpha child (no python-pptx API)."""
    srgb = fore_color._xFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    srgb.append(srgb.makeelement(qn("a:alpha"), {"val": str(int(pct * 1000))}))


def _set_adjust(sp, value: float) -> None:
    """Set a preset shape's first adjustment (e.g. corner radius fraction); no-op if unsupported."""
    try:
        sp.adjustments[0] = value
    except (IndexError, ValueError, ZeroDivisionError):  # pragma: no cover
        pass


def _stadium_adjust(sp) -> None:
    _set_adjust(sp, 0.5)


class _Builder:
    """Holds the slide + the family margin offset, and the per-element emit helpers."""

    def __init__(self, diagram: PositionedDiagram):
        self.d = diagram
        # swimlane/sequence bake absolute coords; generic + ER translate by 24px (matches the SVG).
        self.m = 0.0 if (diagram.lanes or diagram.diagram_type == "sequence") else 24.0
        self.prs = Presentation()
        self.prs.slide_width = Emu(int(round((diagram.width + 2 * self.m) * EMU_PER_PX)))
        self.prs.slide_height = Emu(int(round((diagram.height + 2 * self.m) * EMU_PER_PX)))
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        self.shapes = self.slide.shapes

    def _e(self, px: float) -> Emu:
        return Emu(int(round(px * EMU_PER_PX)))

    def _box(self, x: float, y: float, w: float, h: float) -> tuple[Emu, Emu, Emu, Emu]:
        return (self._e(x + self.m), self._e(y + self.m), self._e(w), self._e(h))

    def _pt(self, p: tuple[float, float]) -> tuple[Emu, Emu]:
        return (self._e(p[0] + self.m), self._e(p[1] + self.m))

    def rect(
        self, shape: MSO_SHAPE, x: float, y: float, w: float, h: float, fill: str | None,
        line: str | None, line_w: float = 1.0, opacity: int | None = None, shadow: bool = False,
    ):
        sp = self.shapes.add_shape(shape, *self._box(x, y, w, h))
        if shadow:
            _add_shadow(sp)
        else:
            sp.shadow.inherit = False
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = _rgb(fill)
            if opacity is not None:
                _set_alpha(sp.fill.fore_color, opacity)
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = _rgb(line)
            sp.line.width = self._e(line_w)
        return sp

    def text_in(self, sp, label, *, size: float = 12.0, color: str = _DEFAULT_TEXT,
                bold: bool = False, align=PP_ALIGN.CENTER, wrap: bool = True) -> None:
        tf = sp.text_frame
        tf.word_wrap = wrap
        for side in ("left", "right", "top", "bottom"):
            setattr(tf, f"margin_{side}", 0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = label if isinstance(label, str) else label.text
        run.font.size = Pt(size)
        run.font.bold = bold
        _set_run_font(run, _FONT)
        run.font.color.rgb = _rgb(color)
        if not isinstance(label, str):
            rtl = _rtl_label(label)
            # tag the run's language so PowerPoint applies the right bidi/shaping (helps mixed
            # Arabic/English spacing); set the paragraph base direction for RTL labels.
            run._r.get_or_add_rPr().set("lang", label.lang or ("ar-SA" if rtl else "en-US"))
            if rtl:
                p._p.get_or_add_pPr().set("rtl", "1")  # no python-pptx API (invariant 4)

    def textbox(self, x: float, y: float, w: float, h: float, label, **kw):
        tb = self.shapes.add_textbox(*self._box(x, y, w, h))
        self.text_in(tb, label, **kw)
        return tb

    def connector(self, p1, p2, color: str, width: float, dashed: bool = False) -> None:
        cn = self.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, *self._pt(p1), *self._pt(p2))
        cn.shadow.inherit = False
        cn.line.color.rgb = _rgb(color)
        cn.line.width = self._e(width)
        if dashed:
            ln = cn.line._get_or_add_ln()
            ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))


# ---------------------------------------------------------------------------
# Element emitters
# ---------------------------------------------------------------------------

def _emit_node(b: _Builder, node: PNode, badge_side: str) -> None:
    if node.shape in ("initial", "final"):
        _emit_pseudostate(b, node)
        return
    sp = b.rect(
        _SHAPE.get(node.shape, MSO_SHAPE.RECTANGLE), node.x, node.y, node.width, node.height,
        _fill_of(node.style), _stroke_of(node.style), _stroke_w_of(node.style),
        shadow=node.shape in _SHADOW_SHAPES,
    )
    if node.shape == "stadium":
        _stadium_adjust(sp)
    if node.shape == "cube":  # label on the front face (matches the engine), not the bbox centre
        d = _CUBE_DEPTH
        b.textbox(node.x, node.y + d, node.width - d, node.height - d, node.label,
                  color=_font_color_of(node.style))
    else:
        b.text_in(sp, node.label, color=_font_color_of(node.style),
                  size=float((node.style.get("text") or {}).get("size", 12)))
    if node.badge:
        _emit_badge(b, node, badge_side)


def _emit_pseudostate(b: _Builder, node: PNode) -> None:
    stroke = _stroke_of(node.style)
    r = min(node.width, node.height) / 2
    cx, cy = node.x + node.width / 2, node.y + node.height / 2
    if node.shape == "initial":
        b.rect(MSO_SHAPE.OVAL, cx - r, cy - r, 2 * r, 2 * r, stroke, None)
        return
    b.rect(MSO_SHAPE.OVAL, cx - r, cy - r, 2 * r, 2 * r, _fill_of(node.style), stroke,
           _stroke_w_of(node.style))
    ir = r * 0.5
    b.rect(MSO_SHAPE.OVAL, cx - ir, cy - ir, 2 * ir, 2 * ir, stroke, None)


def _emit_badge(b: _Builder, node: PNode, side: str) -> None:
    num = (node.badge or "").rstrip(".")
    accent = _stroke_of(node.style)
    cx = node.x + node.width if side == "right" else node.x
    sp = b.rect(MSO_SHAPE.OVAL, cx - _BADGE_R, node.y - _BADGE_R, 2 * _BADGE_R, 2 * _BADGE_R,
                accent, "#FFFFFF", 1.5)
    b.text_in(sp, num, size=11, color="#FFFFFF", bold=True)


def _emit_marker(b: _Builder, m: Marker) -> None:
    if m.kind == "start":
        b.rect(MSO_SHAPE.OVAL, m.cx - m.r, m.cy - m.r, 2 * m.r, 2 * m.r, _MARKER_BLACK, None)
        return
    b.rect(MSO_SHAPE.OVAL, m.cx - m.r, m.cy - m.r, 2 * m.r, 2 * m.r, "#FFFFFF", _MARKER_BLACK, 2.0)
    ir = m.r * 0.45
    b.rect(MSO_SHAPE.OVAL, m.cx - ir, m.cy - ir, 2 * ir, 2 * ir, _MARKER_BLACK, None)


def _emit_swimlane_chrome(b: _Builder, d: PositionedDiagram) -> None:
    rtl = d.direction == "RL"
    vertical = d.orientation == "vertical"
    _emit_title_bar(b, d)
    for group in d.lane_groups:
        sp = b.rect(MSO_SHAPE.ROUNDED_RECTANGLE, group.x, group.y, group.width, group.height,
                    (group.hue or {}).get("label", _PHASE_FILL), None, opacity=92)
        b.text_in(sp, group.label, size=12, color="#FFFFFF", bold=True)
        bodyPr = sp.text_frame._txBody.find(qn("a:bodyPr"))  # always present
        if bodyPr is not None:
            bodyPr.set("vert", "vert270")  # read upward (matches the SVG rotate)
    for band in d.lanes:
        hue = band.hue or {}
        b.rect(MSO_SHAPE.RECTANGLE, band.x, band.y, band.width, band.height,
               hue.get("row", _LANE_ROW), hue.get("label", _LANE_ACCENT), 1.0, opacity=85)
        cx, cy, cw, ch = _chip_rect(band, rtl, vertical)
        chip = b.rect(MSO_SHAPE.ROUNDED_RECTANGLE, cx, cy, cw, ch,
                      hue.get("label", _LANE_ACCENT), None)
        b.text_in(chip, band.label, size=13, color="#FFFFFF", bold=True)
    _emit_separators(b, d, rtl, vertical)


def _chip_rect(band: LaneBand, rtl: bool, vertical: bool) -> tuple[float, float, float, float]:
    if vertical:
        w = band.width - 16.0
        return (band.x + 8.0, band.y + (_V_HEADER - _V_CHIP_H) / 2, w, _V_CHIP_H)
    w = _LABEL_W - 16.0
    x = (band.x + band.width - w - _CHIP_INSET) if rtl else band.x + _CHIP_INSET
    return (x, band.y + (band.height - _CHIP_H) / 2, w, _CHIP_H)


def _emit_title_bar(b: _Builder, d: PositionedDiagram) -> None:
    if not d.title or not d.lanes:
        return
    lanes = d.lanes
    title_x = min([bd.x for bd in lanes] + [g.x for g in d.lane_groups])
    title_right = max(bd.x + bd.width for bd in lanes)
    title_top = d.height - (lanes[-1].y + lanes[-1].height)
    title_bottom = d.phases[0].y if d.phases else lanes[0].y
    title = d.theme.get("title") or {}
    sp = b.rect(MSO_SHAPE.ROUNDED_RECTANGLE, title_x, title_top, title_right - title_x,
                title_bottom - title_top, str(title.get("fill", _TITLE_FILL)), None)
    b.text_in(sp, d.title, size=18, color=str(title.get("text", "#FFFFFF")), bold=True)


def _emit_separators(b: _Builder, d: PositionedDiagram, rtl: bool, vertical: bool) -> None:
    lanes = d.lanes
    m = lanes[0].x
    if vertical:
        sep_y = lanes[0].y + _V_HEADER
        b.connector((lanes[0].x, sep_y), (lanes[-1].x + lanes[-1].width, sep_y), _SEP, 2.0)
        return
    top, bottom = lanes[0].y, lanes[-1].y + lanes[-1].height
    sep_x = (d.width - m - _LABEL_W) if rtl else m + _LABEL_W
    b.connector((sep_x, top), (sep_x, bottom), _SEP, 2.0)
    sep = d.phase_separator or {}
    color = str(sep.get("color", _SEP))
    width = float(sep.get("width", 1.5))
    dashed = sep.get("style") != "solid"
    for phase in d.phases:
        b.connector((phase.x, top), (phase.x, bottom), color, width, dashed)
        _emit_phase(b, phase)
    if d.phases:
        last = max(d.phases, key=lambda p: p.x + p.width)
        edge = last.x + last.width
        b.connector((edge, top), (edge, bottom), color, width, dashed)


def _emit_phase(b: _Builder, phase: PhaseBand) -> None:
    sp = b.rect(MSO_SHAPE.ROUNDED_RECTANGLE, phase.x, phase.y, phase.width, phase.height,
                _PHASE_FILL, None, opacity=92)
    b.text_in(sp, phase.label, size=13, color="#FFFFFF", bold=True)


def _emit_sequence_chrome(b: _Builder, d: PositionedDiagram) -> None:
    if d.title:
        b.textbox(0.0, 0.0, d.width, 24.0, d.title, size=18, color=_DEFAULT_TEXT, bold=True)
    bottom = d.height - 24.0
    for node in d.nodes:
        cx = node.x + node.width / 2
        b.connector((cx, node.y + node.height), (cx, bottom), _SEQ_STEM, 1.5, dashed=True)
    for act in d.activations:
        b.rect(MSO_SHAPE.RECTANGLE, act.x, act.y, act.width, act.height, "#FFFFFF",
               _SEQ_ACT_BORDER, 1.5)


def _emit_entity(b: _Builder, node: PNode) -> None:
    x, y, w, h = node.x, node.y, node.width, node.height
    title_h = node.rows[0].y_offset if node.rows else h
    container = b.rect(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, "#FFFFFF", _ER_BORDER, 1.5)
    # The container's default rounded-rect radius is large; the title's was small -> the title's
    # corners poked past the container's rounding. Pin BOTH to ~6px so the title's rounded top
    # aligns exactly with the container's top. ROUND_2_SAME rounds only the TOP two corners
    # (square bottom) -> rounded-top/square-bottom header like the SVG.
    _set_adjust(container, 6.0 / min(w, h))
    title = b.rect(MSO_SHAPE.ROUND_2_SAME_RECTANGLE, x, y, w, title_h, _ER_TITLE_FILL, None)
    _set_adjust(title, 6.0 / min(w, title_h))
    b.text_in(title, node.label, size=13, color="#FFFFFF", bold=True)
    for r in node.rows:
        ry = y + r.y_offset
        b.connector((x, ry), (x + w, ry), _ER_ROW_SEP, 1.0)
        align = PP_ALIGN.RIGHT if _rtl_label(r.label) else PP_ALIGN.LEFT
        b.textbox(x + 10.0, ry, w - 20.0, r.height, r.label, size=12, color=_DEFAULT_TEXT,
                  align=align)
        if r.key:
            tw = 22.0
            ky = ry + r.height / 2
            pill = b.rect(MSO_SHAPE.ROUNDED_RECTANGLE, x + w - 10.0 - tw, ky - 8, tw, 16,
                          _ER_KEY_FILL.get(r.key, "#777777"), None)
            b.text_in(pill, r.key, size=10, color="#FFFFFF", bold=True)


def _emit_edge(b: _Builder, edge: PositionedEdge, default_w: float, curved: bool) -> None:
    color = str(edge.style.get("stroke", _DEFAULT_EDGE))
    width = float(edge.style.get("width", default_w) or default_w)
    dashed = edge.style.get("style") == "dashed"
    _freeform(b, _rounded_points(list(edge.points), curved), color, width, dashed=dashed)
    if edge.label and edge.label_xy:  # label_xy is already off the line (offset_edge_labels)
        lx, ly = edge.label_xy
        w = max(40.0, len(edge.label.text) * 7.0)  # ~text width; transparent (no white slab)
        b.textbox(lx - w / 2, ly - 8, w, 16.0, edge.label, size=12, color=color, wrap=False)


def _freeform(b: _Builder, pts: list[tuple[float, float]], color: str, width: float,
              *, dashed: bool) -> None:
    if len(pts) < 2:
        return
    off = b.m
    fb = b.shapes.build_freeform(
        (pts[0][0] + off) * EMU_PER_PX, (pts[0][1] + off) * EMU_PER_PX, scale=1.0
    )
    fb.add_line_segments(
        [((px + off) * EMU_PER_PX, (py + off) * EMU_PER_PX) for px, py in pts[1:]], close=False
    )
    sp = fb.convert_to_shape()
    sp.shadow.inherit = False
    sp.fill.background()
    sp.line.color.rgb = _rgb(color)
    sp.line.width = b._e(width)
    ln = sp.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))


def _build(diagram: PositionedDiagram) -> _Prs:
    b = _Builder(diagram)
    if diagram.lanes:
        _emit_swimlane_chrome(b, diagram)
    elif diagram.diagram_type == "sequence":
        _emit_sequence_chrome(b, diagram)
    badge_side = resolve_badge_side(diagram.direction == "RL", diagram.theme)
    default_w = _EDGE_WIDTH_DEFAULT.get(diagram.diagram_type, 1.0)
    curved = resolve_edge_corners(diagram.theme)
    for edge in diagram.edges:  # edges first, shapes on top (arrowheads tuck at borders)
        _emit_edge(b, edge, default_w, curved)
    for node in diagram.nodes:
        if node.rows:
            _emit_entity(b, node)
        else:
            _emit_node(b, node, badge_side)
    for marker in diagram.markers:
        _emit_marker(b, marker)
    return b.prs


def _normalize_zip(raw: bytes) -> bytes:
    """Re-emit the .pptx zip with constant mtimes so the same spec is byte-identical (invariant 7).
    python-pptx stamps each entry with wall-clock time; preserve order + content, fix the rest."""
    src = zipfile.ZipFile(io.BytesIO(raw))
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as zf:
        for name in src.namelist():
            zi = zipfile.ZipInfo(name, date_time=(1980, 1, 1, 0, 0, 0))
            zi.compress_type = zipfile.ZIP_DEFLATED
            zi.external_attr = 0o600 << 16
            zf.writestr(zi, src.read(name))
    return out.getvalue()


def to_pptx_bytes(diagram: PositionedDiagram, meta: dict[str, str] | None = None) -> bytes:
    """Build the .pptx for ``diagram`` and return deterministic bytes (provenance in core props)."""
    prs = _build(diagram)
    cp = prs.core_properties
    cp.title = (meta or {}).get("title") or "Tarseem diagram"
    cp.created = _FIXED_TS
    cp.modified = _FIXED_TS
    if meta:
        cp.comments = "; ".join(f"{k}={v}" for k, v in sorted(meta.items()))
    buf = io.BytesIO()
    prs.save(buf)
    return _normalize_zip(buf.getvalue())


def write_pptx(
    diagram: PositionedDiagram, out_path: str | Path, meta: dict[str, str] | None = None
) -> WriteResult:
    """Write ``diagram`` to ``out_path`` as a native-shape .pptx. Returns the path + a
    CapabilityReport declaring the writer's fidelity ceiling for this diagram."""
    out = Path(out_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_bytes(to_pptx_bytes(diagram, meta))
    return WriteResult(path=out, report=_report(diagram))


def _report(diagram: PositionedDiagram):
    warnings: list[CapabilityWarning] = []
    supports = {
        "shapes": "full",
        "lanes": "full" if diagram.lanes else "none",
        "phases": "full" if diagram.phases else "none",
        "badges": "full",
        "markers": "full",
        "edge_routes": "full",  # exact IR polyline (rounded corners) as a freeform connector
        "edge_labels": "full",
        "curved_edges": "partial",  # rounded corners approximated by sampled segments
        "ports": "partial" if any(n.rows for n in diagram.nodes) else "none",
        "gradients": "none",
        "fonts_embedded": "none",  # names Cairo (a:cs); PowerPoint substitutes if not installed
        "rtl_shaping": "partial",  # pPr rtl=1 + cs font set; shaping delegated to PowerPoint
        "theme_fidelity": "partial",  # flat fills/strokes/text colours; no gradients/tints
        "metadata": "full",  # provenance in core properties
    }
    unknown = sorted({n.shape for n in diagram.nodes
                      if n.shape not in _SHAPE and n.shape not in ("initial", "final", "table")
                      and not getattr(n, "rows", ())})
    for shape in unknown:
        supports["shapes"] = "partial"
        warnings.append(CapabilityWarning("feature-approximated", "shapes",
                                          f"shape {shape!r} drawn as a plain rectangle"))
    if any(n.rows for n in diagram.nodes):
        warnings.append(CapabilityWarning("feature-approximated", "ports",
                        "ER rows are explicit cells; not native table ports"))
    if any(_rtl_label(n.label) for n in diagram.nodes):
        warnings.append(CapabilityWarning("feature-approximated", "rtl_shaping",
                        "a:pPr rtl=1 + a:cs=Cairo set; bidi shaping delegated to PowerPoint"))
    warnings.append(CapabilityWarning("feature-dropped", "fonts_embedded",
                    "PPTX names Cairo; embedding caused PowerPoint repair prompts (reverted)"))
    return build_capability_report("pptx", supports, warnings)
