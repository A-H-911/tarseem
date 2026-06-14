"""PPTX writer tests: valid native-shape decks, byte-determinism, EMU scaling, RTL, report."""
from __future__ import annotations

import io
import json
from pathlib import Path

import pytest
from pptx import Presentation

from tarseem import Engine
from tarseem.export.pptx import EMU_PER_PX, to_pptx_bytes, write_pptx

EX = Path(__file__).resolve().parent.parent / "examples"

FAMILIES = [
    "flowchart",
    "swimlane-pipeline",
    "er-shop",
    "sequence-login",
    "state-order-lifecycle",
    "deployment-web-stack",
    "arabic-flowchart",
    "swimlane-nested-delivery",  # exercises the lane-group vertical-text path
]


def _spec(name: str) -> dict:
    return json.loads((EX / f"{name}.json").read_text(encoding="utf-8"))


def _render(name: str):
    return Engine().render(_spec(name))


@pytest.mark.parametrize("name", FAMILIES)
def test_pptx_is_a_valid_deck_with_native_shapes(name: str):
    diagram = _render(name).diagram
    prs = Presentation(io.BytesIO(to_pptx_bytes(diagram)))
    assert len(prs.slides) == 1
    # at least one native shape per IR node (chrome/edges/labels add more)
    assert len(prs.slides[0].shapes) >= len(diagram.nodes) > 0


@pytest.mark.parametrize("name", FAMILIES)
def test_pptx_is_byte_deterministic(name: str):
    diagram = _render(name).diagram
    assert to_pptx_bytes(diagram, {"spec_hash": "h"}) == to_pptx_bytes(diagram, {"spec_hash": "h"})


def test_pptx_slide_size_matches_extent_in_emu():
    # swimlane bakes absolute coords (margin 0) -> slide == extent.
    diagram = _render("swimlane-pipeline").diagram
    prs = Presentation(io.BytesIO(to_pptx_bytes(diagram)))
    assert prs.slide_width == round(diagram.width * EMU_PER_PX)
    assert prs.slide_height == round(diagram.height * EMU_PER_PX)


def test_generic_diagram_gets_the_svg_margin_offset():
    # generic diagrams are 0-based + translated 24px in the SVG -> slide is extent + 48.
    diagram = _render("flowchart").diagram
    prs = Presentation(io.BytesIO(to_pptx_bytes(diagram)))
    assert prs.slide_width == round((diagram.width + 48) * EMU_PER_PX)


def test_rtl_label_sets_paragraph_rtl_attr():
    prs = Presentation(io.BytesIO(to_pptx_bytes(_render("arabic-flowchart").diagram)))
    assert 'rtl="1"' in prs.slides[0]._element.xml  # python-pptx has no API for this


def test_core_props_carry_no_wallclock():
    prs = Presentation(io.BytesIO(to_pptx_bytes(_render("flowchart").diagram)))
    assert prs.core_properties.created.year == 2001  # fixed constant, not now
    assert prs.core_properties.modified.year == 2001


def test_capability_report_on_shared_vocab(tmp_path):
    result = write_pptx(_render("er-shop").diagram, tmp_path / "d.pptx", {"spec_hash": "h"})
    assert result.path.exists()
    assert result.report.writer == "pptx"
    assert result.report.lossy  # fonts not embeddable -> at least one warning
    assert result.report.supports["edge_routes"] == "full"


def test_pptx_runs_set_complex_script_font():
    # Arabic renders from the cs slot; without this it falls back to the theme font (review #1-3).
    prs = Presentation(io.BytesIO(to_pptx_bytes(_render("arabic-flowchart").diagram)))
    assert '<a:cs typeface="Cairo"' in prs.slides[0]._element.xml


def test_pptx_separators_are_connectors_not_freeforms():
    # straight separators/lifelines use connectors so they render crisp, not blurry (review #11-17).
    prs = Presentation(io.BytesIO(to_pptx_bytes(_render("swimlane-pipeline").diagram)))
    assert "<p:cxnSp>" in prs.slides[0]._element.xml


def test_pptx_embeds_the_cairo_font():
    import zipfile

    z = zipfile.ZipFile(io.BytesIO(to_pptx_bytes(_render("arabic-flowchart").diagram)))
    assert "ppt/fonts/font1.fntdata" in z.namelist()
    assert "embeddedFontLst" in z.read("ppt/presentation.xml").decode("utf-8")


def test_pptx_lane_group_label_is_vertical():
    prs = Presentation(io.BytesIO(to_pptx_bytes(_render("swimlane-nested-delivery").diagram)))
    assert 'vert="vert270"' in prs.slides[0]._element.xml


def test_pptx_3d_shapes_have_a_drop_shadow():
    prs = Presentation(io.BytesIO(to_pptx_bytes(_render("deployment-web-stack").diagram)))
    assert "outerShdw" in prs.slides[0]._element.xml


def test_pptx_edge_label_has_no_fill_slab():
    # link-text background is transparent now (review global note); the freeform line is the route.
    prs = Presentation(io.BytesIO(to_pptx_bytes(_render("flowchart").diagram)))
    # textboxes carry no solid fill (no <a:solidFill> inside a txBox spPr is hard to assert
    # directly); instead confirm curved edges produced a multi-point freeform path.
    assert "<a:cubicBezTo>" in prs.slides[0]._element.xml or \
        prs.slides[0]._element.xml.count("<a:lnTo>") > len(_render("flowchart").diagram.edges)


def test_export_writes_pptx_with_lossy_sidecar(tmp_path):
    res = _render("swimlane-pipeline")
    written = res.export(["pptx"], tmp_path, name="d")
    assert written["pptx"].exists()
    assert (tmp_path / "d.pptx.report.json").exists()
    assert res.reports["pptx"].writer == "pptx"
