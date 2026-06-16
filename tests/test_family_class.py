"""UML class family (Sub-stage 6): validate -> compile -> measure -> render."""
from __future__ import annotations

import shutil

import pytest

from tarseem import Engine
from tarseem.measure import measure_graph
from tarseem.model.compile import compile_spec
from tarseem.validation import validate

SPEC = {
    "specVersion": "1.0",
    "diagramType": "class",
    "nodes": [
        {
            "id": "User",
            "label": {"text": "User"},
            "attributes": ["- id: int", "- name: String"],
            "methods": ["+ getName(): String", "+ setName(n: String)"],
        },
        {"id": "Account", "label": {"text": "Account"}, "attributes": ["- balance: float"]},
    ],
    "edges": [{"source": "User", "target": "Account", "label": {"text": "owns"}}],
}

requires_node = pytest.mark.skipif(shutil.which("node") is None, reason="ELK layout needs Node")


def test_validate_accepts_class_strings_and_methods():
    assert validate(SPEC).ok


def test_compile_parses_members_not_er_rows():
    g = compile_spec(SPEC)
    user = next(n for n in g.nodes if n.id == "User")
    assert user.rows == ()  # class nodes do not become ER tables
    assert [m.group for m in user.members] == ["attr", "attr", "method", "method"]
    assert user.members[0].label.text == "- id: int"
    assert user.members[2].label.text == "+ getName(): String"


def test_measure_stamps_increasing_member_geometry():
    g = measure_graph(compile_spec(SPEC))
    user = next(n for n in g.nodes if n.id == "User")
    assert user.width and user.height
    ys = [m.y_offset for m in user.members]
    assert ys == sorted(ys) and ys[0] > 0  # below the name bar, monotonically increasing
    assert all(m.height > 0 for m in user.members)


def test_class_box_not_uniform_resized():
    # class boxes are content-sized; uniformNodeSize must not snap them (like ER tables).
    spec = {**SPEC, "layout": {"uniformNodeSize": True}}
    g = measure_graph(compile_spec(spec))
    widths = {n.id: n.width for n in g.nodes}
    assert widths["User"] != widths["Account"]  # different content -> different size


@requires_node
def test_render_class_svg_has_compartments_and_dividers():
    svg = Engine().render(SPEC).svg
    assert "User" in svg and "- id: int" in svg and "+ getName(): String" in svg
    # name|attributes|methods dividers (>=2 lines for the User box) plus the edge
    assert svg.count("<line") >= 2
    assert "owns" in svg


@requires_node
def test_class_drawio_emits_explicit_compartment_cells():
    from tarseem.export.drawio import to_drawio_xml

    xml = to_drawio_xml(Engine().render(SPEC).diagram)
    assert "class_User" in xml  # container
    assert "classtitle_User" in xml  # grey name bar
    assert "classmember_User_" in xml  # member text cells
    assert "classdiv_User_" in xml  # compartment dividers
    assert "+ getName(): String" in xml


@requires_node
def test_class_pptx_is_a_valid_deck():
    import io

    from pptx import Presentation

    from tarseem.export.pptx import to_pptx_bytes

    prs = Presentation(io.BytesIO(to_pptx_bytes(Engine().render(SPEC).diagram)))
    assert len(prs.slides) == 1
    assert "getName" in prs.slides[0]._element.xml
