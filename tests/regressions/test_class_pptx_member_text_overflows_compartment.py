"""Regression: UML class members overlapped / overflowed their compartment in PPTX.

Root cause (export layer only — the IR/measure/layout are correct; the SVG renders fine):
each member textbox was emitted with word-wrap ON (``wrap="square"``) and the python-pptx
``add_textbox`` default ``<a:spAutoFit/>`` (resize-shape-to-fit-text). When PowerPoint's font
metrics wrapped a long member to two lines, spAutoFit grew the textbox taller, pushing it into
the next member's row → visible overlap, text not staying inside the compartment.

Fix: member textboxes stay a single line (``wrap=False``) and do not auto-grow
(``auto_size = MSO_AUTO_SIZE.NONE`` → ``<a:noAutofit/>``), so the box keeps its stamped
geometry and members never overlap. (Horizontal spill when Cairo is not installed is the
documented PPTX fonts ceiling, not this bug.)
"""
from __future__ import annotations

import io

import pytest
from pptx import Presentation
from pptx.oxml.ns import qn

from tarseem import Engine
from tarseem.export.pptx import to_pptx_bytes

pytestmark = pytest.mark.skipif(
    __import__("shutil").which("node") is None, reason="ELK layout needs Node"
)

# Minimal repro: a class with a long method that PowerPoint would wrap.
SPEC = {
    "specVersion": "1.0",
    "diagramType": "class",
    "nodes": [
        {
            "id": "Order",
            "label": {"text": "Order"},
            "attributes": ["- id: int", "- total: Money"],
            "methods": ["+ addLine(p: Product): void", "+ checkout()"],
        }
    ],
    "edges": [],
}

_MEMBER_TEXTS = {
    "- id: int", "- total: Money", "+ addLine(p: Product): void", "+ checkout()",
}


def _member_textboxes(prs):
    for sh in prs.slides[0].shapes:
        if sh.has_text_frame and sh.text_frame.text in _MEMBER_TEXTS:
            yield sh


def test_class_pptx_member_textboxes_wrap_and_shrink_to_stay_inside():
    prs = Presentation(io.BytesIO(to_pptx_bytes(Engine().render(SPEC).diagram)))
    boxes = list(_member_textboxes(prs))
    assert len(boxes) == 4, "expected one textbox per class member"
    for sh in boxes:
        body = sh.text_frame._txBody.find(qn("a:bodyPr"))
        assert body is not None
        # wrap ON: a long member wraps at the box edge instead of spilling past the right border
        assert body.get("wrap") == "square", f"{sh.text_frame.text!r} should word-wrap"
        autofit = [c.tag.split("}")[-1] for c in body]
        # shrink-to-fit, never grow: normAutofit shrinks the text; spAutoFit (grow box) is the bug
        assert "spAutoFit" not in autofit, f"{sh.text_frame.text!r} must not grow the box"
        assert "normAutofit" in autofit, f"{sh.text_frame.text!r} should shrink-to-fit"


def test_class_pptx_member_rows_keep_stamped_non_overlapping_geometry():
    # the boxes themselves must stay at the IR-stamped, non-overlapping y positions.
    diagram = Engine().render(SPEC).diagram
    prs = Presentation(io.BytesIO(to_pptx_bytes(diagram)))
    tops = sorted(sh.top for sh in _member_textboxes(prs))
    # each member box is below the previous one by at least its height (no vertical overlap)
    node = diagram.nodes[0]
    row_h_emu = round(node.members[0].height * 9525)
    for a, b in zip(tops, tops[1:], strict=False):
        assert b - a >= row_h_emu - 1, "member textboxes must not overlap vertically"


def test_class_member_font_is_px_scaled_to_points_not_oversized():
    # The layout is CSS px @ 96 dpi; the PPTX run font size must be px * 72/96 pt. Setting Pt(px)
    # directly renders every label 1.33x too large (16 px for a 12 px size) → members overflow
    # their box. This is the root cause of the "text floating outside the class" report.
    from pptx.util import Pt

    prs = Presentation(io.BytesIO(to_pptx_bytes(Engine().render(SPEC).diagram)))
    box = next(_member_textboxes(prs))
    run = box.text_frame.paragraphs[0].runs[0]
    assert run.font.size == Pt(12 * 72 / 96), "12 px member must render at 9 pt, not 12 pt"


def test_class_member_boxes_have_horizontal_headroom_for_wider_renderers():
    # The box must be wider than the bundled-Cairo text so a viewer that renders the
    # (non-embedded) font slightly wider — e.g. PowerPoint with the installed Cairo — keeps the
    # left-aligned member inside its right border (the "text floating outside the class" bug).
    from tarseem.measure import _shared_measurer

    node = Engine().render(SPEC).diagram.nodes[0]
    measurer = _shared_measurer()
    inner = node.width - 24.0  # geometry CLASS_PAD_X (12) each side
    widest = max(measurer.width(m.label.text, 12.0) for m in node.members)
    assert inner >= widest * 1.10, "class box lacks headroom for wider non-embedded fonts"
