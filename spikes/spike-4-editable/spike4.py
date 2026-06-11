"""Spike 4: editable-export probes (throwaway).

(A) Hand-built draw.io file: native swimlane POOL + 2 lanes + an RTL Arabic node
    (uncompressed mxGraphModel, writingDirection=rtl). Rendered headlessly through the
    real diagrams.net GraphViewer engine to prove it opens correctly.
(B) Minimal python-pptx deck: 2 shapes + a connector + an RTL paragraph patched with
    <a:pPr rtl="1"> via lxml (python-pptx has no API for it). Re-opened to assert the
    rtl flag survives in the saved XML.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Inches
from playwright.sync_api import sync_playwright

HERE = Path(__file__).resolve().parent
OUT = HERE / "out"
OUT.mkdir(exist_ok=True)

VIEWER = "https://viewer.diagrams.net/js/viewer-static.min.js"


# ---- (A) draw.io native swimlane pool ---------------------------------------
def build_drawio() -> str:
    # horizontal=0 -> horizontal swimlane (vertical title bar on the left, lanes as rows)
    pool = 'swimlane;html=1;horizontal=0;startSize=24;rounded=0;fillColor=none;swimlaneFillColor=#ffffff;'
    lane = 'swimlane;html=1;horizontal=0;startSize=24;writingDirection=rtl;'
    node = 'rounded=1;whiteSpace=wrap;html=1;writingDirection=rtl;fillColor=#C8F0E0;strokeColor=#269973;'
    model = f'''<mxGraphModel dx="800" dy="500" grid="0" guides="1" tooltips="1" connect="1" arrows="1" fold="1" page="1" pageScale="1" pageWidth="720" pageHeight="320" math="0" shadow="0">
  <root>
    <mxCell id="0"/>
    <mxCell id="1" parent="0"/>
    <mxCell id="pool" value="العملية" style="{pool}" vertex="1" parent="1">
      <mxGeometry x="40" y="40" width="640" height="240" as="geometry"/>
    </mxCell>
    <mxCell id="lane1" value="مقدّم الطلب" style="{lane}" vertex="1" parent="pool">
      <mxGeometry x="24" y="0" width="616" height="120" as="geometry"/>
    </mxCell>
    <mxCell id="lane2" value="المراجعة" style="{lane}" vertex="1" parent="pool">
      <mxGeometry x="24" y="120" width="616" height="120" as="geometry"/>
    </mxCell>
    <mxCell id="task1" value="تعبئة الطلب" style="{node}" vertex="1" parent="lane1">
      <mxGeometry x="80" y="35" width="150" height="50" as="geometry"/>
    </mxCell>
    <mxCell id="task2" value="مراجعة" style="{node}" vertex="1" parent="lane2">
      <mxGeometry x="380" y="35" width="150" height="50" as="geometry"/>
    </mxCell>
    <mxCell id="e1" style="edgeStyle=orthogonalEdgeStyle;rounded=1;html=1;endArrow=classic;strokeColor=#269973;" edge="1" parent="1" source="task1" target="task2">
      <mxGeometry relative="1" as="geometry"/>
    </mxCell>
  </root>
</mxGraphModel>'''
    return model


def validate_drawio(model: str) -> dict:
    """Deterministic structural proof the file is a valid native-swimlane RTL drawio."""
    import xml.etree.ElementTree as ET

    root = ET.fromstring(model)  # raises if not well-formed
    cells = root.findall(".//mxCell")
    by_id = {c.get("id"): c for c in cells}

    def style(cid: str) -> str:
        return by_id[cid].get("style") or ""

    lanes = [c for c in cells if "swimlane" in (c.get("style") or "") and c.get("parent") == "pool"]
    return {
        "well_formed_xml": True,
        "pool_is_swimlane": "swimlane" in style("pool"),
        "two_lanes": len(lanes) == 2,
        "lanes_horizontal0": all("horizontal=0" in (c.get("style") or "") for c in lanes),
        "rtl_on_lane": any("writingDirection=rtl" in (c.get("style") or "") for c in lanes),
        "node_parented_to_lane": by_id["task1"].get("parent") == "lane1",
        "rtl_on_node": "writingDirection=rtl" in style("task1"),
        "edge_present": any(c.get("edge") == "1" for c in cells),
    }


def write_drawio(model: str) -> Path:
    mxfile = f'<mxfile host="tarseem-spike-4">\n  <diagram name="Pool" id="spike4">\n{model}\n  </diagram>\n</mxfile>\n'
    path = OUT / "pool.drawio"
    path.write_text(mxfile, encoding="utf-8")
    return path


def render_drawio(model: str) -> bool:
    """Render via the real diagrams.net GraphViewer; screenshot proves it opens correctly."""
    import json

    cfg = json.dumps({"highlight": "#0000ff", "nav": False, "resize": True, "xml": model})
    html = (
        "<!doctype html><html><head><meta charset='utf-8'>"
        f"<script src='{VIEWER}'></script></head><body style='margin:0'>"
        f"<div class='mxgraph' style='max-width:100%;border:1px solid #eee;' data-mxgraph='{cfg}'></div>"
        "</body></html>"
    )
    import functools
    import http.server
    import socketserver
    import threading

    (OUT / "pool_viewer.html").write_text(html, encoding="utf-8")
    handler = functools.partial(http.server.SimpleHTTPRequestHandler, directory=str(OUT))
    httpd = socketserver.TCPServer(("127.0.0.1", 0), handler)
    port = httpd.server_address[1]
    threading.Thread(target=httpd.serve_forever, daemon=True).start()
    try:
        with sync_playwright() as p:
            b = p.chromium.launch()
            pg = b.new_page(device_scale_factor=2)
            pg.goto(f"http://127.0.0.1:{port}/pool_viewer.html")
            pg.wait_for_selector(".mxgraph svg", timeout=8000)
            pg.wait_for_timeout(800)
            pg.query_selector(".mxgraph").screenshot(path=str(OUT / "pool_drawio.png"))
            b.close()
        return True
    except Exception as exc:  # network/viewer failure must not be silent
        print(f"  [drawio viewer render FAILED: {exc}]")
        return False
    finally:
        httpd.shutdown()


# ---- (B) python-pptx deck with RTL paragraph --------------------------------
def set_rtl(shape) -> None:
    for para in shape.text_frame.paragraphs:
        pPr = para._p.get_or_add_pPr()
        pPr.set("rtl", "1")
        pPr.set("algn", "r")
        for run in para.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.append(rPr.makeelement(qn("a:cs"), {"typeface": "Arial"}))


def build_pptx() -> Path:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    s1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(2.4), Inches(1))
    s1.text_frame.text = "مقدّم الطلب"
    s2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(1.5), Inches(2.4), Inches(1))
    s2.text_frame.text = "المراجعة"
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(3.4), Inches(2.0), Inches(5.5), Inches(2.0))
    conn.begin_connect(s1, 3)
    conn.end_connect(s2, 1)
    for sh in (s1, s2):
        set_rtl(sh)
    path = OUT / "deck.pptx"
    prs.save(path)
    return path


def verify_pptx(path: Path) -> dict:
    prs = Presentation(str(path))
    slide = prs.slides[0]
    xml = slide._element.xml
    shapes = list(slide.shapes)
    n_conn = sum(1 for sh in shapes if sh.element.tag.endswith("}cxnSp"))
    return {
        "rtl_flag_present": 'rtl="1"' in xml,
        "algn_r_present": 'algn="r"' in xml,
        "cs_typeface_present": "<a:cs" in xml,
        "n_shapes": len(shapes),
        "n_connectors": n_conn,
    }


def main() -> int:
    model = build_drawio()
    dpath = write_drawio(model)
    print(f"drawio written: {dpath.name} ({dpath.stat().st_size} bytes, uncompressed)")
    dchecks = validate_drawio(model)
    print("drawio structural assertions:")
    for k, val in dchecks.items():
        print(f"  {k}: {val}")
    rendered = render_drawio(model)
    print(f"drawio GraphViewer visual render (best-effort): {rendered} -> out/pool_drawio.png")

    ppath = build_pptx()
    print(f"pptx written: {ppath.name} ({ppath.stat().st_size} bytes)")
    v = verify_pptx(ppath)
    print("pptx re-open assertions:")
    for k, val in v.items():
        print(f"  {k}: {val}")

    drawio_ok = all(dchecks.values())
    pptx_ok = v["rtl_flag_present"] and v["algn_r_present"] and v["n_connectors"] >= 1
    ok = drawio_ok and pptx_ok
    print(f"\ndrawio structural: {drawio_ok} | pptx: {pptx_ok} | visual render: {rendered}")
    print(f"spike-4 automated checks pass: {ok}  (visual drawio open in diagrams.net = manual step)")
    return 0 if ok else 1


if __name__ == "__main__":
    raise SystemExit(main())
